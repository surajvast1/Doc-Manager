import io
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

def extract_text_from_s3(file_stream, file_extension):
    if file_extension == '.csv':
        df = pd.read_csv(file_stream)
        return df.to_string(index=False)
    elif file_extension in ['.xlsx', '.xls']:
        df = pd.read_excel(file_stream)
        return df.to_string(index=False)
    elif file_extension == '.pdf':
        pdf_reader = PdfReader(file_stream)
        return " ".join(page.extract_text() for page in pdf_reader.pages)
    elif file_extension == '.pptx':
        presentation = Presentation(file_stream)
        return " ".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
    elif file_extension == '.docx':
        doc = Document(file_stream)
        return " ".join(para.text for para in doc.paragraphs)
    else:
        return None
